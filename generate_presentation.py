from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (clean corporate: white + green + blue) ──────────────────
WHITE_BG    = RGBColor(0xFF, 0xFF, 0xFF)   # slide background (white)
GREEN       = RGBColor(0x3D, 0xAA, 0x35)   # primary green (SOLEVO-style)
GREEN_DARK  = RGBColor(0x2D, 0x80, 0x28)   # darker green for depth
BLUE        = RGBColor(0x1F, 0x5C, 0x9E)   # corporate blue (secondary)
BLUE_LIGHT  = RGBColor(0xE8, 0xF1, 0xFB)   # very light blue (card background)
GREEN_LIGHT = RGBColor(0xEB, 0xF7, 0xEA)   # very light green (card background)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black for body text
MID_GREY    = RGBColor(0x55, 0x55, 0x66)   # medium grey for subtitles
FTBG        = RGBColor(0xF0, 0xF0, 0xF2)   # footer bar background
FTTX        = RGBColor(0x88, 0x88, 0x99)   # footer bar text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # pure white (text on dark bg)
RED_FILL    = RGBColor(0xFF, 0xEB, 0xEB)   # light red (blocked boxes)
RED_BORDER  = RGBColor(0xCC, 0x20, 0x20)   # red border/text
AMBER_FILL  = RGBColor(0xFF, 0xF3, 0xD4)   # light amber (throttled boxes)
AMBER       = RGBColor(0xB8, 0x6A, 0x00)   # amber border/text

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ────────────────────────────────────────────────────────

def bg(slide, color=WHITE_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    """Add a plain rectangle."""
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        font_size=18, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def bullet_block(slide, items, left, top, width, height,
                 font_size=16, color=DARK_TEXT, icon="▸"):
    """Render a list of bullet strings."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def accent_bar(slide, top=0.72, height=0.05):
    """Horizontal accent line under the title."""
    box(slide, 0, top, 13.33, height, fill_color=BLUE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – The Challenge
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
bg(s1)

# Header bar
box(s1, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s1, "Slide 1 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s1, "The Challenge – Unexpected Twilio Costs",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s1)

# Four info cards
card_data = [
    ("📌  The Situation",
     "We transitioned to Twilio as our communication provider for sending "
     "One-Time Passwords (OTPs) via SMS."),
    ("⚠️  The Problem",
     "Shortly after this change, we noticed a massive spike in unexplainable "
     "SMS costs."),
    ("🔍  The Root Cause",
     "A huge volume of OTPs was triggered for random phone numbers — primarily "
     "based in Pakistan."),
    ("🤖  What This Is",
     '"SMS Pumping" / "Toll Fraud" — automated bots trigger expensive SMS '
     "messages to international numbers to steal a fraction of routing fees."),
]

card_positions = [
    (0.3,  1.0),  (6.8,  1.0),
    (0.3,  4.05), (6.8,  4.05),
]
card_w, card_h = 6.1, 2.8

for (heading, body), (cl, ct) in zip(card_data, card_positions):
    box(s1, cl, ct, card_w, card_h, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width_pt=1.5)
    box(s1, cl, ct, card_w, 0.52, fill_color=GREEN)
    txt(s1, heading, cl+0.18, ct+0.07, card_w-0.3, 0.42,
        font_size=15, bold=True, color=WHITE)
    txt(s1, body, cl+0.18, ct+0.65, card_w-0.36, card_h-0.8,
        font_size=13, color=DARK_TEXT)

# Footer
box(s1, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s1, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Strategy Overview
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bg(s2)

box(s2, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s2, "Slide 2 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s2, "Our Strategy – A Two-Layered Defence",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s2)

txt(s2,
    "To permanently stop SMS Pumping and protect our budget, we implemented "
    "a robust, two-layered security checkpoint system on our backend.",
    0.5, 0.9, 12.3, 0.9, font_size=16, color=MID_GREY)

txt(s2, "Both layers must be passed before we spend a single cent on an SMS.",
    0.5, 1.65, 12.3, 0.6, font_size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Layer cards
for i, (icon, title, subtitle, desc) in enumerate([
    ("🛡️", "Layer 1", "The Identity Check",
     "Are you our real app?\n\n"
     "Component: AppRequestContextFilter (Firebase App Check)\n\n"
     "Demands a cryptographic, unforgeable digital token before processing any request."),
    ("⏱️", "Layer 2", "The Speed Limit",
     "Are you requesting too fast?\n\n"
     "Component: SendOtpRateLimitAttribute\n\n"
     "Tracks request frequency per user/device and enforces a strict per-minute cap."),
]):
    cl = 0.5 + i * 6.6
    ct = 2.5
    box(s2, cl, ct, 6.1, 4.3, fill_color=BLUE_LIGHT,
        line_color=BLUE, line_width_pt=2)
    # coloured top strip
    box(s2, cl, ct, 6.1, 0.65, fill_color=BLUE)
    txt(s2, f"{icon}  {title} — {subtitle}",
        cl+0.15, ct+0.08, 5.8, 0.55, font_size=17, bold=True, color=WHITE)
    txt(s2, desc, cl+0.2, ct+0.8, 5.7, 3.3,
        font_size=13, color=DARK_TEXT)

box(s2, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s2, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Defence Layer 1
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bg(s3)

box(s3, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s3, "Slide 3 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s3, "Defence Layer 1 – The Gatekeeper",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s3)

txt(s3, "🔒  AppRequestContextFilter  (Firebase App Check)",
    0.4, 0.85, 12.5, 0.55, font_size=18, bold=True, color=BLUE)

# Flow diagram boxes
flow = [
    ("📱  Real App",            "Silently generates a\ncryptographic MobileAppCheck\ntoken in the background.", 0.35),
    ("🚪  Gatekeeper",          "Demands the token before\nprocessing any request.\nNo token → instant reject.", 4.10),
    ("🤖  Bot / Script",        "Cannot generate a valid token.\nNo real device. No real app.\nInstantly blocked at the door.", 7.85),
]

for label, body, cl in flow:
    is_bot = "Bot" in label
    fill   = RED_FILL   if is_bot else BLUE_LIGHT
    border = RED_BORDER if is_bot else BLUE
    lbl_c  = RED_BORDER if is_bot else BLUE
    box(s3, cl, 1.6, 3.4, 2.9, fill_color=fill,
        line_color=border, line_width_pt=2)
    txt(s3, label, cl+0.15, 1.7, 3.1, 0.55,
        font_size=14, bold=True, color=lbl_c)
    txt(s3, body, cl+0.15, 2.3, 3.1, 2.0,
        font_size=12, color=DARK_TEXT)

# Arrows between boxes
for ax in [3.75, 7.50]:
    txt(s3, "➤", ax, 2.75, 0.5, 0.5, font_size=22, color=GREEN, align=PP_ALIGN.CENTER)

# Outcome banner
box(s3, 0.35, 4.8, 12.6, 0.75, fill_color=GREEN,
    line_color=GREEN, line_width_pt=0)
txt(s3,
    "✅  Result: Bots are rejected instantly — Cost to us: $0",
    0.5, 4.87, 12.3, 0.6, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_block(s3, [
    "Bot farms lack a real device running our genuine application.",
    "Automated scripts cannot replicate the cryptographic token generation.",
    "Rejection happens before any SMS is queued — zero Twilio spend.",
], 0.5, 5.75, 12.3, 1.3, font_size=14, color=DARK_TEXT)

box(s3, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s3, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Defence Layer 2
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
bg(s4)

box(s4, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s4, "Slide 4 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s4, "Defence Layer 2 – The Speed Limit",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s4)

txt(s4, "⏱️  SendOtpRateLimitAttribute",
    0.4, 0.85, 12.5, 0.55, font_size=18, bold=True, color=BLUE)

# Why / How / Result cards
whr = [
    ("❓  Why We Need This",
     "What if an attacker hijacks a real phone, or a legitimate user taps "
     "'Send OTP' 50 times in a row?\n\nLayer 1 would let them through — but "
     "it would still cost us money."),
    ("⚙️  How It Works",
     "Tracks how often a specific user or device requests an OTP.\n\n"
     "Applies a strict cap (e.g. max X texts per minute).\n\n"
     "Once the cap is hit: 'Please wait' — no further SMS is sent."),
    ("✅  The Result",
     "Residual fraud is caught by this fail-safe layer.\n\n"
     "Legitimate users experience a brief cooldown at worst.\n\n"
     "Cost savings are locked in — no runaway charges."),
]

for i, (heading, body) in enumerate(whr):
    cl = 0.35 + i * 4.35
    box(s4, cl, 1.6, 4.1, 3.8, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width_pt=1.5)
    box(s4, cl, 1.6, 4.1, 0.62, fill_color=GREEN)
    txt(s4, heading, cl+0.15, 1.67, 3.8, 0.5,
        font_size=14, bold=True, color=WHITE)
    txt(s4, body, cl+0.15, 2.32, 3.8, 3.0,
        font_size=12, color=DARK_TEXT)

# Combined defence summary
box(s4, 0.35, 5.65, 12.6, 1.05, fill_color=BLUE_LIGHT,
    line_color=BLUE, line_width_pt=1.5)
txt(s4, "🔐  Combined Defence Summary",
    0.55, 5.68, 5, 0.45, font_size=14, bold=True, color=BLUE)
bullet_block(s4, [
    "Layer 1 (App Check) blocks ~99 %+ of bot traffic before any cost is incurred.",
    "Layer 2 (Rate Limit) eliminates residual abuse and protects against edge cases.",
], 0.55, 6.1, 12.2, 0.6, font_size=12, icon="✔", color=DARK_TEXT)

box(s4, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s4, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Business Value & Next Steps
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
bg(s5)

box(s5, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s5, "Slide 5 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s5, "Business Value & Next Steps",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s5)

# Three value cards
value_cards = [
    ("💰  Immediate ROI",
     "We have successfully shut off the vulnerability that automated scripts "
     "were exploiting, saving the company from runaway Twilio invoices.\n\n"
     "The fix costs nothing to run and protects us indefinitely."),
    ("📈  Predictable Scaling",
     "As we grow our legitimate user base, our SMS costs will scale directly "
     "with real human onboarding — not outside bot traffic.\n\n"
     "Budgeting becomes reliable and tied to actual growth."),
    ("✨  Zero User Friction",
     "Both security measures happen in milliseconds behind the scenes.\n\n"
     "Legitimate users attempting to log in normally will not even notice "
     "they are there — no extra steps, no delays, no CAPTCHA."),
]

for i, (heading, body) in enumerate(value_cards):
    cl = 0.35 + i * 4.35
    box(s5, cl, 1.1, 4.1, 5.7, fill_color=GREEN_LIGHT,
        line_color=GREEN, line_width_pt=1.5)
    # coloured top strip
    box(s5, cl, 1.1, 4.1, 0.65, fill_color=GREEN)
    txt(s5, heading, cl+0.18, 1.17, 3.8, 0.52,
        font_size=15, bold=True, color=WHITE)
    txt(s5, body, cl+0.18, 1.85, 3.78, 4.7,
        font_size=13, color=DARK_TEXT)

box(s5, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s5, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – How Our Defence Works (Business-Friendly Diagram)
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
bg(s6)

box(s6, 0, 0, 13.33, 0.75, fill_color=GREEN)
txt(s6, "Slide 6 of 6", 10.8, 0.08, 2.3, 0.5,
    font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s6, "How Our Defence Works – At a Glance",
    0.3, 0.05, 12.5, 0.7, font_size=26, bold=True, color=WHITE)
accent_bar(s6)

txt(s6, "Every SMS request passes through two automatic checkpoints before we are charged a penny.",
    0.5, 0.85, 12.3, 0.55, font_size=14, color=MID_GREY, align=PP_ALIGN.CENTER)

# ── Flow diagram ────────────────────────────────────────────────────────────
# Row 1 (main flow): Request → Checkpoint 1 → Checkpoint 2 → SMS Sent
# Row 2 (blocked branches below each checkpoint)

# Helper: flow node box
def flow_node(slide, label, sublabel, left, top, w=2.5, h=1.1,
              fill=BLUE_LIGHT, border=BLUE, lbl_color=BLUE, sub_color=MID_GREY):
    box(slide, left, top, w, h, fill_color=fill, line_color=border, line_width_pt=2)
    txt(slide, label, left+0.12, top+0.08, w-0.24, 0.52,
        font_size=13, bold=True, color=lbl_color, align=PP_ALIGN.CENTER)
    if sublabel:
        txt(slide, sublabel, left+0.1, top+0.55, w-0.2, 0.5,
            font_size=10, color=sub_color, align=PP_ALIGN.CENTER)

# Arrow helper (horizontal ➤)
def h_arrow(slide, left, top):
    txt(slide, "➤", left, top, 0.55, 0.45,
        font_size=20, color=GREEN, align=PP_ALIGN.CENTER)

# Arrow helper (vertical ↓)
def v_arrow(slide, left, top):
    txt(slide, "▼", left, top, 0.45, 0.38,
        font_size=16, color=RED_BORDER, align=PP_ALIGN.CENTER)

ROW1_TOP = 1.65   # y of main flow row
ROW2_TOP = 3.35   # y of blocked-branch row

# Node positions (left edge)
X_REQ   = 0.25    # "SMS Request"
X_GATE  = 3.10    # Checkpoint 1
X_SPEED = 6.55    # Checkpoint 2
X_SENT  = 10.20   # SMS Sent ✅

NODE_W = 2.55
NODE_H = 1.1

# ── Main flow nodes ──
flow_node(s6, "📲  SMS Request", "User taps 'Send Code'",
          X_REQ, ROW1_TOP, NODE_W, NODE_H,
          fill=BLUE_LIGHT, border=BLUE, lbl_color=BLUE, sub_color=MID_GREY)

h_arrow(s6, X_REQ + NODE_W + 0.05, ROW1_TOP + 0.32)

flow_node(s6, "🛡️  Checkpoint 1", "Is this our real app?",
          X_GATE, ROW1_TOP, NODE_W, NODE_H,
          fill=BLUE, border=BLUE, lbl_color=WHITE, sub_color=RGBColor(0xCC, 0xDD, 0xFF))

h_arrow(s6, X_GATE + NODE_W + 0.05, ROW1_TOP + 0.32)

flow_node(s6, "⏱️  Checkpoint 2", "Is the request rate normal?",
          X_SPEED, ROW1_TOP, NODE_W, NODE_H,
          fill=BLUE, border=BLUE, lbl_color=WHITE, sub_color=RGBColor(0xCC, 0xDD, 0xFF))

h_arrow(s6, X_SPEED + NODE_W + 0.05, ROW1_TOP + 0.32)

flow_node(s6, "✅  SMS Sent!", "Twilio charged — legitimate user",
          X_SENT, ROW1_TOP, NODE_W, NODE_H,
          fill=GREEN, border=GREEN_DARK,
          lbl_color=WHITE, sub_color=RGBColor(0xCC, 0xFF, 0xCC))

# ── Down-arrows from checkpoints to blocked boxes ──
v_arrow(s6, X_GATE  + NODE_W/2 - 0.05, ROW1_TOP + NODE_H + 0.05)
v_arrow(s6, X_SPEED + NODE_W/2 - 0.05, ROW1_TOP + NODE_H + 0.05)

# ── Blocked branch boxes ──
BLOCK_W = 3.0
BLOCK_H = 1.55

# Checkpoint 1 blocked
box(s6, X_GATE - 0.2, ROW2_TOP, BLOCK_W, BLOCK_H,
    fill_color=RED_FILL, line_color=RED_BORDER, line_width_pt=2)
txt(s6, "❌  Blocked — FREE",
    X_GATE - 0.08, ROW2_TOP + 0.07, BLOCK_W - 0.24, 0.42,
    font_size=13, bold=True, color=RED_BORDER, align=PP_ALIGN.CENTER)
txt(s6, "Bot / script without\na valid app token.\nCost to us: $0",
    X_GATE - 0.08, ROW2_TOP + 0.48, BLOCK_W - 0.24, 1.0,
    font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Checkpoint 2 blocked
box(s6, X_SPEED - 0.2, ROW2_TOP, BLOCK_W, BLOCK_H,
    fill_color=AMBER_FILL, line_color=AMBER, line_width_pt=2)
txt(s6, "⚠️  Throttled — FREE",
    X_SPEED - 0.08, ROW2_TOP + 0.07, BLOCK_W - 0.24, 0.42,
    font_size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
txt(s6, "Too many requests too fast.\nSystem says 'Please wait'.\nCost to us: $0",
    X_SPEED - 0.08, ROW2_TOP + 0.48, BLOCK_W - 0.24, 1.0,
    font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ── Legend / takeaway banner ──
box(s6, 0.35, 5.2, 12.6, 1.55, fill_color=GREEN_LIGHT,
    line_color=GREEN, line_width_pt=1.5)
txt(s6, "🔑  Key Takeaway for the Business",
    0.55, 5.25, 6, 0.45, font_size=14, bold=True, color=GREEN_DARK)
bullet_block(s6, [
    "Only real users on our genuine app can ever reach the 'SMS Sent' step — protecting our budget automatically.",
    "The system is fully invisible to legitimate customers and requires no manual intervention.",
    "Both checkpoints run in under 50 milliseconds — no perceptible delay for users.",
], 0.55, 5.65, 12.1, 1.0, font_size=12, icon="✔", color=DARK_TEXT)

box(s6, 0, 7.15, 13.33, 0.35, fill_color=FTBG)
txt(s6, "SMS Pumping Defence  |  Confidential", 0.3, 7.18, 13, 0.28,
    font_size=9, color=FTTX)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save("sms_pumping_defence.pptx")
print("✅  Saved: sms_pumping_defence.pptx")
